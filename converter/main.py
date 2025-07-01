from typing import List
from pptx import Presentation
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
import glob
import os
import sys

SLIDE_WIDTH = 75


def parse(file_path: str) -> List[str]:
    prs = Presentation(file_path)
    print(f"Processing: {file_path} (slides 1-{len(prs.slides)})...")
    print('-' * 50)
    
    doc = []
    divider = f"{'-' * SLIDE_WIDTH}\n{'-' * SLIDE_WIDTH}\n"
    for slide in prs.slides:
        slide_text = ""
        for el in slide.shapes:
            if hasattr(el, "text"):
                slide_text += f"{el.text}\n"
        if len(slide_text) > 0:
            doc.append(slide_text)

    return doc


def format(file_path: str, text: List[str], dir_path: str):
    doc_name = file_path[:-4].split('/')[-1]
    document = Document()
    document.add_heading(doc_name, level=0)
    document.add_heading("Savannah PRE Worksheet", level=1)
    weird_char_count = 0
    for slide in text:
        lines = [l for l in slide.split("\n") if l]
        for line in lines:
            try:
                if ':' in line:
                    before, after = line.split(':', 1)
                    para = document.add_paragraph()
                    run = para.add_run(before)
                    run.bold = True
                    para.add_run(f":{after}")
                else:
                    document.add_paragraph(line)
            except:
                weird_char_count += 1
        
        lb = document.add_paragraph()
        add_horizontal_line(lb)
    
    if weird_char_count > 0:
        print(f"Warning: found {weird_char_count} weird characters")
    
    dir_path_processed = f"{dir_path}_processed"
    if not os.path.exists(dir_path_processed):
        os.mkdir(dir_path_processed)
    doc_path = f"{dir_path_processed}/{doc_name}docx"
    document.save(doc_path)
    
    
def add_horizontal_line(paragraph):
    p = paragraph._p  # Access the paragraph's XML element
    pPr = p.get_or_add_pPr()  # Get or create the paragraph properties
    pbdr = OxmlElement('w:pBdr')  # Create the paragraph border element
    bottom = OxmlElement('w:bottom')  # Create the bottom border element

    # Set attributes for the bottom border (this creates the horizontal line)
    bottom.set(qn('w:val'), 'single')  # Single line border
    bottom.set(qn('w:sz'), '6')  # Thickness (in eighths of a point)
    bottom.set(qn('w:space'), '1')  # Spacing between text and border
    bottom.set(qn('w:color'), 'auto')  # Automatic color (black)

    pbdr.append(bottom)  # Add the bottom border to the paragraph borders
    pPr.append(pbdr)  # Add the paragraph borders to the paragraph properties


def find_pptx_files(dir_path):
    pptx_files = []
    for root, dirs, files in os.walk(dir_path):
        for file in files:
            if file.endswith('.pptx'):# or file.endswith('.ppt'):
                pptx_files.append(os.path.join(root, file))
    return pptx_files


def process_batch(dir_path: str):
    files = find_pptx_files(dir_path)
    for file in files:
        text = parse(file)
        format(file, text, dir_path)


if __name__ == "__main__":
    args = sys.argv
    if len(args) <= 1:
        print("Usage: python3 main.py dir_name\nSearches all subdirectories for pptx files, saving the output to the directory `dir_name_processed`")
    else:
        for dir_name in args[1:]:
            print(f"Processing: {dir_name}")
            process_batch("crime_and_punishment")
        print(f"Saving results to: {dir_name}_processed/")
